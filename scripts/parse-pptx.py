"""
AutoPresent Studio — PPTX 파서
python-pptx를 사용하여 PPTX의 슬라이드별 콘텐츠를 구조적으로 추출합니다.
사용법: python parse-pptx.py <input.pptx> [output_dir]
"""

import json
import os
import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


# 27종 AutoPresent 슬라이드 타입 매핑 규칙
TYPE_MAPPING = {
    "title": "cover",
    "toc": "toc",
    "bullets_only": "bullets",
    "numbers_heavy": "stats",
    "table_dominant": "table",
    "chart_dominant": "stats",
    "image_dominant": "hero",
    "comparison": "problem-solution",
    "quote": "quote",
    "team": "team",
    "timeline": "timeline",
    "closing": "conclusion",
}


def parse_pptx(pptx_path, output_dir="extracted"):
    """PPTX 파일을 파싱하여 content-brief.json 형식으로 반환"""
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)

    sections = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "title": None,
            "bodyTexts": [],
            "bulletLists": [],
            "tables": [],
            "images": [],
            "speakerNote": None,
        }

        # 발표자 노트 추출
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                slide_data["speakerNote"] = notes_frame.text.strip()

        # 각 도형 분석
        for shape in slide.shapes:
            _parse_shape(shape, slide_num, output_dir, slide_data)

        # 우세 콘텐츠 타입 판별
        dominant = _detect_dominant(slide_data)

        # 섹션 구성
        content_text = "\n".join(slide_data["bodyTexts"])
        key_points = []
        for bl in slide_data["bulletLists"]:
            key_points.extend(bl)

        section = {
            "sectionIndex": slide_num,
            "heading": slide_data["title"] or f"슬라이드 {slide_num}",
            "content": content_text or slide_data["speakerNote"] or "",
            "keyPoints": key_points,
            "dataPoints": [],
            "images": slide_data["images"],
            "suggestedSlideType": _map_type(dominant, slide_num, len(prs.slides)),
        }

        # 표 데이터를 dataPoints로 변환
        for table in slide_data["tables"]:
            for row in table.get("rows", []):
                if len(row) >= 2:
                    section["dataPoints"].append({
                        "label": str(row[0]),
                        "value": str(row[1]),
                        "source": "PPTX 추출",
                    })

        sections.append(section)

    # content-brief.json 형식
    brief = {
        "sourceType": "pptx",
        "title": sections[0]["heading"] if sections else "제목 없음",
        "summary": f"PPTX에서 추출한 {len(sections)}개 슬라이드 콘텐츠",
        "sections": sections,
        "metadata": {
            "totalPages": len(prs.slides),
            "language": "ko",
            "extractedImages": sum(len(s["images"]) for s in sections),
            "processingNotes": [f"python-pptx로 파싱, {len(prs.slides)}장 처리"],
        },
    }

    return brief


def _parse_shape(shape, slide_num, output_dir, slide_data):
    """개별 도형 파싱"""
    # 제목 프레임
    if shape.has_text_frame:
        tf = shape.text_frame
        text = tf.text.strip()
        if not text:
            return

        # 제목 판별 (placeholder 타입 또는 큰 폰트)
        is_title = False
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            ph_type = shape.placeholder_format.type
            if ph_type in (1, 15, 13):  # TITLE, SUBTITLE, CENTER_TITLE
                is_title = True

        if is_title and not slide_data["title"]:
            slide_data["title"] = text
        else:
            # 불릿 리스트 감지
            paragraphs = [p.text.strip() for p in tf.paragraphs if p.text.strip()]
            if len(paragraphs) >= 3:
                has_bullets = any(
                    p.level > 0 or p.text.startswith(("•", "-", "–", "※"))
                    for p in tf.paragraphs
                    if p.text.strip()
                )
                if has_bullets:
                    slide_data["bulletLists"].append(paragraphs)
                else:
                    slide_data["bodyTexts"].append(text)
            else:
                slide_data["bodyTexts"].append(text)

    # 표
    elif shape.has_table:
        table = shape.table
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        rows = []
        for row in list(table.rows)[1:]:
            rows.append([cell.text.strip() for cell in row.cells])
        slide_data["tables"].append({"headers": headers, "rows": rows})

    # 이미지
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            filename = f"slide{slide_num:03d}_img_{shape.shape_id}.{ext}"
            filepath = os.path.join(output_dir, filename)
            with open(filepath, "wb") as f:
                f.write(image.blob)
            slide_data["images"].append({
                "path": filepath,
                "caption": "",
                "ocrText": "",
                "confidence": 0.0,
            })
        except Exception:
            pass


def _detect_dominant(slide_data):
    """슬라이드의 우세 콘텐츠 타입 판별"""
    if slide_data["tables"]:
        return "table_dominant"
    if slide_data["images"] and len(slide_data["images"]) >= 2:
        return "image_dominant"
    if slide_data["bulletLists"]:
        return "bullets_only"

    # 숫자 비율 확인
    all_text = " ".join(slide_data["bodyTexts"])
    numbers = re.findall(r"\d+[%만억원$]", all_text)
    if len(numbers) >= 3:
        return "numbers_heavy"

    return "bullets_only"


def _map_type(dominant, slide_num, total_slides):
    """AutoPresent 27종 타입 매핑"""
    if slide_num == 1:
        return "cover"
    if slide_num == 2:
        return "toc"
    if slide_num == total_slides:
        return "conclusion"
    return TYPE_MAPPING.get(dominant, "bullets")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python parse-pptx.py <input.pptx> [output_dir]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "extracted"

    result = parse_pptx(pptx_path, output_dir)

    # content-brief.json 저장
    output_path = os.path.join(os.path.dirname(pptx_path), "content-brief.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"✅ PPTX 파싱 완료: {len(result['sections'])}개 슬라이드")
    print(f"📁 결과: {output_path}")
    print(f"🖼️  추출 이미지: {result['metadata']['extractedImages']}개")
